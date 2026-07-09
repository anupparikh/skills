#!/usr/bin/env python3
"""lint_brand.py — brand linter for PowerPoint decks built by deck.py.

Checks a rendered .pptx against its source template + brand tokens: off-palette
colors, off-brand fonts, footer/page-number chrome, canvas bounds, shape
collisions, text overflow, dead space (large under-filled slides), and logo
placement. Pure XML/geometry inspection — no LLM calls, no rendering. Permanent
verifier asset; run it after every deck.py change or content build.

    python3 lint_brand.py DECK.pptx --template TEMPLATE.pptx [--config brand.json]
                           [--logo logo.png] [--require-logo-on all|noncover|cover|none]
                           [--no-deadspace] [--json] [--selfcheck]

Exit 0 = clean, 1 = defects found, 2 = usage/IO error.
"""
import argparse, hashlib, json, os, re, sys, zipfile
from collections import namedtuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

EMU_IN = 914400
CANVAS_W, CANVAS_H = 13.333, 7.5     # inches, 16:9 — deck.py's target canvas
CANVAS_TOL = 0.03                    # inches, float-rounding slack
FOOTER_BAND_TOP = 6.6                # inches; shape top below this = footer/page-number band
OVERFLOW_TOL = 1.08                  # flag only clear overflow: >8% past available height
COLLISION_AREA_FRAC = 0.05           # partial-overlap defect threshold (fraction of smaller bbox)
COLLISION_MIN_PEN = 0.08             # inches; overlaps shallower than this are touching seams, not collisions
                                     # (declared boxes abut with a hair of slack; real glyphs don't reach the edge)
LOGO_POS_TOL = 0.1                   # inches
DEADSPACE_TOP, DEADSPACE_BOTTOM = 0.5, 6.9   # inches; content region, excludes title/logo band above and footer band below
DEADSPACE_GRID_W, DEADSPACE_GRID_H = 40, 24  # coarse rasterization for largest-empty-rectangle sweep (~0.33in x 0.27in cells)
DEADSPACE_FILL_MAX = 0.55            # flag when covered fraction of the content region falls below this.
                                     # Calibrated against ground truth (2026-07-05): the APPROVED human
                                     # decks (source examples/) ship intentional airy slides down to ~54%
                                     # fill — figure-right/bottom-breathing evidence slides. At 0.60 the
                                     # check flagged the human decks' own good slides (15% of the
                                     # case-studies deck). 0.55 flags only slides below the human floor
                                     # (the genuinely-broken 26-48% cases) while accepting the airy-but-
                                     # intentional look the approved design language uses. Do NOT raise
                                     # back to 0.60 to make a deck "cleaner" — that flags the ground truth.
DEADSPACE_CARD_AREA_CAP = 0.35       # non-text filled shape counts as a "card" only below this frac of content-region area;
                                     # bigger = a full-bleed background rect, not real content
DEADSPACE_TEXT_FORWARD_MAX = 0.30    # relaxed fill threshold for text-forward slides (a native body
                                     # placeholder carrying real prose — bullets/agenda archetypes).
                                     # Their ink-area fill is structurally low (text lines, not tiles),
                                     # so the 0.55 gate made them unwinnable (v0.1 field report #7);
                                     # 0.30 still catches a 2-bullet slide rattling in an empty frame.
FONT_DIRS = ["/System/Library/Fonts", "/System/Library/Fonts/Supplemental",
             "/Library/Fonts", os.path.expanduser("~/Library/Fonts")]
DEFAULT_PCS_TEMPLATE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "..", "..", "clients", "pcs", "brand", "precision-cell-systems-deck-template.pptx")

Defect = namedtuple("Defect", "slide check detail shape")


def in_(emu):
    """EMU -> inches"""
    return emu / EMU_IN


def I(v):
    """inches -> EMU int"""
    return int(v * EMU_IN)


# ---- allowed-set derivation (template + brand config are the source of truth) -----

def _read_zip_texts(path, pattern):
    z = zipfile.ZipFile(path)
    texts = [z.read(n).decode("utf-8", "ignore") for n in z.namelist() if re.match(pattern, n)]
    z.close()
    return texts


def theme_colors(template_path):
    x = _read_zip_texts(template_path, r"ppt/theme/theme1\.xml$")[0]
    return {m.upper() for m in re.findall(r'<a:srgbClr val="([0-9A-Fa-f]{6})"', x)}


def theme_fonts(template_path):
    x = _read_zip_texts(template_path, r"ppt/theme/theme1\.xml$")[0]
    mj = re.search(r'<a:majorFont><a:latin typeface="([^"]*)"', x)
    mn = re.search(r'<a:minorFont><a:latin typeface="([^"]*)"', x)
    return {f.group(1) for f in (mj, mn) if f}


def masters_layouts_colors(template_path):
    hexes = set()
    for x in _read_zip_texts(template_path, r"ppt/slide(Master|Layout)\d*\.xml$"):
        hexes |= {m.upper() for m in re.findall(r'<a:srgbClr val="([0-9A-Fa-f]{6})"', x)}
    return hexes


def masters_layouts_fonts(template_path):
    fonts = set()
    for x in _read_zip_texts(template_path, r"ppt/slide(Master|Layout)\d*\.xml$"):
        # a:latin typeface only — a:buFont (bullet-glyph font, e.g. Arial) is not a text typeface
        fonts |= set(re.findall(r'<a:latin typeface="([^"]*)"', x))
    return {f for f in fonts if not f.startswith("+")}   # +mj-lt/+mn-lt are theme refs, not literal names


def config_colors(config_path):
    if not config_path:
        return set()
    pal = json.load(open(config_path)).get("presentation", {}).get("palette", {})
    return {v.lstrip("#").upper() for v in pal.values()
            if isinstance(v, str) and re.fullmatch(r"#?[0-9A-Fa-f]{6}", v)}


def config_fonts(config_path):
    if not config_path:
        return set()
    fonts = json.load(open(config_path)).get("presentation", {}).get("fonts", {})
    return {v for v in fonts.values() if isinstance(v, str)}


def build_allowed(template_path, config_path):
    colors = (theme_colors(template_path) | masters_layouts_colors(template_path)
              | config_colors(config_path) | {"FFFFFF", "000000"})
    fonts = theme_fonts(template_path) | masters_layouts_fonts(template_path) | config_fonts(config_path)
    return colors, fonts


# ---- shape traversal helpers -------------------------------------------------

def _stype(sh):
    try:
        return sh.shape_type
    except Exception:
        return None


def iter_shapes(shapes):
    """Yield every shape, recursing into groups (deck.py itself never groups, but
    hand-built decks might)."""
    for sh in shapes:
        yield sh
        if _stype(sh) == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(sh.shapes)


def _has_text(sh):
    return bool(getattr(sh, "has_text_frame", False)) and sh.text_frame.text.strip()


def _is_candidate(sh):
    """collision-check candidates: text frames with text, and pictures."""
    return _stype(sh) == MSO_SHAPE_TYPE.PICTURE or _has_text(sh)


def _bbox_in(sh):
    return (in_(sh.left), in_(sh.top), in_(sh.left) + in_(sh.width), in_(sh.top) + in_(sh.height))


# ---- checks -------------------------------------------------------------------

def check_canvas_size(prs):
    w, h = in_(prs.slide_width), in_(prs.slide_height)
    if abs(w - CANVAS_W) > CANVAS_TOL or abs(h - CANVAS_H) > CANVAS_TOL:
        return [Defect(None, "canvas", f"deck canvas {w:.3f}x{h:.3f}in != {CANVAS_W}x{CANVAS_H}in", None)]
    return []


def check_palette(prs, allowed):
    defects, seen = [], set()
    for i, slide in enumerate(prs.slides, 1):
        for sh in iter_shapes(slide.shapes):
            for el in sh._element.findall(".//" + qn("a:srgbClr")):
                hexv = (el.get("val") or "").upper()
                if hexv and hexv not in allowed and (i, sh.name, hexv) not in seen:
                    seen.add((i, sh.name, hexv))
                    defects.append(Defect(i, "palette", f"off-palette color #{hexv}", sh.name))
    return defects


def check_fonts(prs, allowed):
    defects = []
    for i, slide in enumerate(prs.slides, 1):
        for sh in iter_shapes(slide.shapes):
            if not getattr(sh, "has_text_frame", False):
                continue
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    name = r.font.name
                    if name and name not in allowed:
                        snippet = (r.text or "")[:30]
                        defects.append(Defect(i, "fonts", f"off-brand font {name!r} on text {snippet!r}", sh.name))
    return defects


def check_footer(prs):
    defects = []
    for i, slide in enumerate(prs.slides, 1):
        band = [sh for sh in iter_shapes(slide.shapes)
                 if sh.top is not None and in_(sh.top) > FOOTER_BAND_TOP - 1e-6]
        page_nums, nonempty = [], False
        for sh in band:
            if not getattr(sh, "has_text_frame", False):
                continue
            if sh.text_frame.text.strip():
                nonempty = True
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    t = (r.text or "").strip()
                    if re.fullmatch(r"\d+", t):
                        page_nums.append(int(t))
        if i == 1:
            if page_nums or nonempty:
                defects.append(Defect(i, "footer", "cover slide has footer/page-number chrome", None))
        else:
            if not page_nums:
                defects.append(Defect(i, "footer", "missing page number in bottom band", None))
            if not nonempty:
                defects.append(Defect(i, "footer", "missing footer text in bottom band", None))
            for pn in page_nums:
                if pn != i:
                    defects.append(Defect(i, "footer", f"page number {pn} != slide index {i}", None))
    return defects


def check_canvas(prs):
    """Top-level shape bboxes only — group-child coordinates live in the group's own
    child coordinate space and need a transform to compare against the slide canvas;
    out of scope here since deck.py never groups shapes."""
    defects = []
    W, H = in_(prs.slide_width), in_(prs.slide_height)
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if None in (sh.left, sh.top, sh.width, sh.height):
                continue
            l, t, w, h = in_(sh.left), in_(sh.top), in_(sh.width), in_(sh.height)
            if l < -CANVAS_TOL or t < -CANVAS_TOL or l + w > W + CANVAS_TOL or t + h > H + CANVAS_TOL:
                defects.append(Defect(i, "canvas",
                    f"shape bbox ({l:.2f},{t:.2f})-({l+w:.2f},{t+h:.2f}) exceeds canvas {W:.3f}x{H:.3f}in",
                    sh.name))
    return defects


def check_collision(prs):
    defects = []
    for i, slide in enumerate(prs.slides, 1):
        cands = [sh for sh in iter_shapes(slide.shapes)
                 if None not in (sh.left, sh.top, sh.width, sh.height) and _is_candidate(sh)]
        for ai in range(len(cands)):
            for bi in range(ai + 1, len(cands)):
                a, b = cands[ai], cands[bi]
                ax0, ay0, ax1, ay1 = _bbox_in(a)
                bx0, by0, bx1, by1 = _bbox_in(b)
                iw = max(0.0, min(ax1, bx1) - max(ax0, bx0))
                ih = max(0.0, min(ay1, by1) - max(ay0, by0))
                inter = iw * ih
                if inter <= 0:
                    continue
                if min(iw, ih) < COLLISION_MIN_PEN:
                    continue   # shallow abutment, not a real overlap
                a_area, b_area = (ax1 - ax0) * (ay1 - ay0), (bx1 - bx0) * (by1 - by0)
                smaller = min(a_area, b_area)
                if smaller <= 0 or inter >= smaller - 1e-6:
                    continue   # full containment is fine (text on a card is the design)
                frac = inter / smaller
                if frac > COLLISION_AREA_FRAC:
                    defects.append(Defect(i, "collision",
                        f"{a.name!r} overlaps {b.name!r} ({frac:.0%} of smaller shape)",
                        f"{a.name} / {b.name}"))
    return defects


# ---- overflow: text-height estimation -----------------------------------------

_FONT_FILE_CACHE = {}


def find_font_file(family, bold=False, italic=False):
    """Best-effort lookup of a .ttf/.ttc/.otf for `family` under the system font dirs,
    by filename match (not a real font-name index — good enough for estimation)."""
    if not family:
        return None
    key = (family, bold, italic)
    if key in _FONT_FILE_CACHE:
        return _FONT_FILE_CACHE[key]
    style = " ".join(s for s, on in (("Bold", bold), ("Italic", italic)) if on)
    wanted = [f"{family} {style}".strip(), family]
    found = None
    for d in FONT_DIRS:
        if not os.path.isdir(d):
            continue
        try:
            files = os.listdir(d)
        except OSError:
            continue
        for w in wanted:
            for fn in files:
                if os.path.splitext(fn)[0].lower() == w.lower() and fn.lower().endswith((".ttf", ".ttc", ".otf")):
                    found = os.path.join(d, fn)
                    break
            if found:
                break
        if found:
            break
        for fn in files:  # loose fallback: filename starts with the family name
            if fn.lower().startswith(family.lower()) and fn.lower().endswith((".ttf", ".ttc", ".otf")):
                found = os.path.join(d, fn)
                break
        if found:
            break
    _FONT_FILE_CACHE[key] = found
    return found


_PIL_FONT_CACHE = {}


def _pil_font(family, pt_size, bold, italic):
    key = (family, round(pt_size), bold, italic)
    if key in _PIL_FONT_CACHE:
        return _PIL_FONT_CACHE[key]
    font = None
    path = find_font_file(family, bold, italic)
    if path:
        try:
            from PIL import ImageFont
            font = ImageFont.truetype(path, size=max(1, int(round(pt_size))))
        except Exception:
            font = None
    _PIL_FONT_CACHE[key] = font
    return font


def _text_width_in(text, family, pt_size, bold=False, italic=False):
    """Rendered width of `text` in inches at `pt_size` pt. Real font metrics via PIL
    when a matching .ttf is found; else avg-char-width heuristic (0.5 * pt_size)."""
    font = _pil_font(family, pt_size, bold, italic)
    if font is not None:
        try:
            bbox = font.getbbox(text)
            return (bbox[2] - bbox[0]) / 72.0
        except Exception:
            pass
    return len(text) * pt_size * 0.5 / 72.0


def estimate_text_height(text_frame, avail_w_in, default_font="Franklin Gothic Book"):
    """Estimate the rendered height (inches) of `text_frame` wrapped to `avail_w_in`
    inches of width. Wraps each paragraph word-by-word using real font metrics (PIL +
    a located system .ttf) when available, else the avg-char-width heuristic. Line
    height ~= 1.2x point size (explicit paragraph line-spacing overrides this), and
    normAutofit fontScale is applied to shrink effective sizes first. This is a
    conservative estimator, not a layout engine — seeds future measured-layout work
    (e.g. swap in an actual text-shaping library once precision matters more than speed).
    """
    scale = 1.0
    bodyPr = text_frame._txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        af = bodyPr.find(qn("a:normAutofit"))
        if af is not None and af.get("fontScale"):
            scale = int(af.get("fontScale")) / 100000.0
    wrap = getattr(text_frame, "word_wrap", True)
    w_in = avail_w_in if wrap is not False else float("inf")

    total = 0.0
    for p in text_frame.paragraphs:
        runs = p.runs
        words = []
        for r in runs:
            sz = (r.font.size.pt if r.font.size else 18) * scale
            fam = r.font.name or default_font
            bold = bool(r.font.bold)
            for wtxt in (r.text or "").split(" "):
                if wtxt:
                    words.append((wtxt, sz, fam, bold))
        if not words:
            total += (18 * scale * 1.2) / 72.0   # empty paragraph still takes a line
            continue
        max_sz, line_w, n_lines = words[0][1], 0.0, 1
        for wtxt, sz, fam, bold in words:
            ww = _text_width_in(wtxt, fam, sz, bold)
            space = _text_width_in(" ", fam, sz, bold) if line_w > 0 else 0.0
            if line_w > 0 and line_w + space + ww > w_in:
                n_lines += 1
                line_w = ww
            else:
                line_w += space + ww
            max_sz = max(max_sz, sz)
        line_h_in = (1.2 * max_sz) / 72.0
        ls = p.line_spacing
        if isinstance(ls, (int, float)):
            line_h_in = (ls * max_sz) / 72.0
        para_h = n_lines * line_h_in
        if p.space_before:
            para_h += p.space_before.pt / 72.0
        if p.space_after:
            para_h += p.space_after.pt / 72.0
        total += para_h
    return total


def check_overflow(prs):
    defects = []
    for i, slide in enumerate(prs.slides, 1):
        for sh in iter_shapes(slide.shapes):
            if not getattr(sh, "has_text_frame", False) or not sh.text_frame.text.strip():
                continue
            if None in (sh.width, sh.height):
                continue
            tf = sh.text_frame
            w_in, h_in = in_(sh.width), in_(sh.height)
            lm = in_(tf.margin_left) if tf.margin_left is not None else 0.1
            rm = in_(tf.margin_right) if tf.margin_right is not None else 0.1
            tm = in_(tf.margin_top) if tf.margin_top is not None else 0.05
            bm = in_(tf.margin_bottom) if tf.margin_bottom is not None else 0.05
            avail_w = max(0.05, w_in - lm - rm)
            avail_h = max(0.05, h_in - tm - bm)
            est = estimate_text_height(tf, avail_w)
            if est > OVERFLOW_TOL * avail_h:
                defects.append(Defect(i, "overflow",
                    f"estimated text height {est:.2f}in > {OVERFLOW_TOL}x available {avail_h:.2f}in",
                    sh.name))
    return defects


# ---- deadspace: large unshaped empty regions -----------------------------------
# The design-critic experiment found "dead space" (title at top, most of the slide
# blank below; a card sized 3x its content) to be the #1 recurring defect, and this
# linter had no way to see it — every other check looks for something wrong that's
# THERE (overflow, collision, off-palette), never for content that's missing. This
# check estimates how much of the content region (excludes the title/logo band and
# footer band) is actually filled, and where the biggest empty block sits.
#
# A shape's "filled" footprint is not always its declared box: a text box is
# routinely drawn much taller than the line or two of text it holds (the "oversized
# card" failure), so a naive bbox-union reads that padding as ink. We estimate the
# rendered text height (reusing estimate_text_height from the overflow check) and
# position it inside the box per its vertical anchor, so an oversized box with short
# text only "fills" the part it actually draws into. Non-text shapes with a visible
# fill (a colored card background, an icon chip) still count at their full declared
# box — a colored rectangle occupies its area regardless of any text inside it —
# except shapes covering most of the content region, which are full-bleed background
# panels, not content, and are excluded (see DEADSPACE_CARD_AREA_CAP).
#
# Calibrated against 7 real decks against ground-truth "dead-space" critic tags
# (see MASTER_PLAN / handoff notes for the experiment): DEADSPACE_FILL_MAX=0.60 gave
# zero false positives on ~90 slides across narrative/dense/singulator x rigid/free
# cells, catching every clearly-empty slide the critic named while staying quiet on
# balanced, densely-packed ones. A largest-empty-rectangle threshold (as an AND
# condition) was tried first per the original spec, but real decks showed the
# defect is often *diffuse* under-fill (several oversized cards, no single big empty
# block) rather than one contiguous void — so fill ratio alone is the gate; the
# empty rectangle is still computed and reported for a useful "which band" detail.


def _ink_bbox(sh):
    """Effective filled bbox for a text-bearing shape: full box width, but height is
    the estimated rendered-text height (not the declared box height), positioned per
    the text frame's vertical anchor (top/center/bottom; default top)."""
    x0, y0, x1, y1 = _bbox_in(sh)
    tf = sh.text_frame
    w_in, h_in = x1 - x0, y1 - y0
    lm = in_(tf.margin_left) if tf.margin_left is not None else 0.1
    rm = in_(tf.margin_right) if tf.margin_right is not None else 0.1
    tm = in_(tf.margin_top) if tf.margin_top is not None else 0.05
    bm = in_(tf.margin_bottom) if tf.margin_bottom is not None else 0.05
    avail_w = max(0.05, w_in - lm - rm)
    est_h = min(h_in, estimate_text_height(tf, avail_w) + tm + bm)
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    anchor = bodyPr.get("anchor") if bodyPr is not None else None
    if anchor == "ctr":
        pad = (h_in - est_h) / 2
        ny0, ny1 = y0 + pad, y0 + pad + est_h
    elif anchor == "b":
        ny0, ny1 = y1 - est_h, y1
    else:
        ny0, ny1 = y0, y0 + est_h
    return x0, ny0, x1, ny1


def _has_visible_fill(sh):
    try:
        return sh.fill.type in (MSO_FILL_TYPE.SOLID, MSO_FILL_TYPE.GRADIENT,
                                 MSO_FILL_TYPE.PATTERNED, MSO_FILL_TYPE.PICTURE)
    except Exception:
        return False


def _boxes_intersect(a, b):
    return not (a[2] <= b[0] or b[2] <= a[0] or a[3] <= b[1] or b[3] <= a[1])


def _deadspace_bboxes(slide, region_area):
    """Content bboxes clipped to the deadspace content band, in the same (x0,y0,x1,y1)
    inches shape as _bbox_in. Text shapes use ink height; pictures use the full box.

    A non-text filled rectangle (e.g. a `stat` tile's separate background shape) earns
    full-box credit ONLY when real text ink sits on it — otherwise it is an empty or
    decorative tile and gets none. Without this, an empty filler box counts as
    "content" and a visibly blank slide passes the fill check (Codex finding, 2026-07-03).

    A text-bearing shape that is ITSELF filled (deck.py's compose `card` is one rounded
    rectangle containing its text) is a deliberate on-brand panel: it earns its full
    declared box, not just its text band — ink-only credit read slides full of visible
    cards as "empty" (v0.1 field report #9). The full-bleed-background area cap still
    applies."""
    text_boxes, pics, fills = [], [], []
    for sh in iter_shapes(slide.shapes):
        if None in (sh.left, sh.top, sh.width, sh.height):
            continue
        is_pic = _stype(sh) == MSO_SHAPE_TYPE.PICTURE
        has_txt = getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip()
        if has_txt:
            bb = _bbox_in(sh)
            if (_has_visible_fill(sh)
                    and (bb[2] - bb[0]) * (bb[3] - bb[1]) <= DEADSPACE_CARD_AREA_CAP * region_area):
                text_boxes.append(bb)          # filled card with text: full-box credit
            else:
                text_boxes.append(_ink_bbox(sh))
        elif is_pic:
            pics.append(_bbox_in(sh))
        elif _has_visible_fill(sh):
            bb = _bbox_in(sh)
            if (bb[2] - bb[0]) * (bb[3] - bb[1]) <= DEADSPACE_CARD_AREA_CAP * region_area:
                fills.append(bb)   # candidate — kept below only if it backs real text
    kept_fills = [f for f in fills if any(_boxes_intersect(f, t) for t in text_boxes)]
    boxes = []
    for x0, y0, x1, y1 in text_boxes + pics + kept_fills:
        y0c, y1c = max(y0, DEADSPACE_TOP), min(y1, DEADSPACE_BOTTOM)
        if y1c > y0c:
            boxes.append((x0, y0c, x1, y1c))
    return boxes


def _deadspace_grid_metrics(boxes, canvas_w):
    """Rasterize the content region to a DEADSPACE_GRID_W x DEADSPACE_GRID_H boolean
    grid of covered cells; return (fill_ratio, largest_empty_rect_frac, rect_cells).
    A coarse grid approximation of a maximal-rectangle-in-histogram sweep — good
    enough to locate "which band is empty" without a real layout engine."""
    gw, gh = DEADSPACE_GRID_W, DEADSPACE_GRID_H
    region_h = DEADSPACE_BOTTOM - DEADSPACE_TOP
    cw, ch = canvas_w / gw, region_h / gh
    covered = [[False] * gw for _ in range(gh)]
    for x0, y0, x1, y1 in boxes:
        c0 = max(0, int(x0 / cw)); c1 = min(gw, int((x1 - 1e-9) / cw) + 1)
        r0 = max(0, int((y0 - DEADSPACE_TOP) / ch)); r1 = min(gh, int((y1 - DEADSPACE_TOP - 1e-9) / ch) + 1)
        for r in range(r0, r1):
            row = covered[r]
            for c in range(c0, c1):
                row[c] = True
    filled = sum(row.count(True) for row in covered)
    fill_ratio = filled / (gw * gh)

    heights = [0] * gw
    best, best_rc = 0, None
    for r in range(gh):
        for c in range(gw):
            heights[c] = 0 if covered[r][c] else heights[c] + 1
        stack = []
        for c in range(gw + 1):
            h = heights[c] if c < gw else 0
            start = c
            while stack and stack[-1][1] >= h:
                sc, sh_ = stack.pop()
                area = sh_ * (c - sc)
                if area > best:
                    best, best_rc = area, (sc, c, r - sh_ + 1, r + 1)
                start = sc
            stack.append((start, h))
    empty_frac = best / (gw * gh)
    return fill_ratio, empty_frac, best_rc


def _band_label(rc, gw, gh):
    """Rough location name (thirds grid) for a (col0,col1,row0,row1) empty rectangle."""
    c0, c1, r0, r1 = rc
    cx, cy = (c0 + c1) / 2 / gw, (r0 + r1) / 2 / gh
    vband = "top" if cy < 1 / 3 else ("bottom" if cy > 2 / 3 else "middle")
    hband = "left" if cx < 1 / 3 else ("right" if cx > 2 / 3 else "center")
    if vband == "middle" and hband == "center":
        return "center"
    if hband == "center":
        return vband
    if vband == "middle":
        return hband
    return f"{vband}-{hband}"


def _is_text_forward(slide):
    """True when a native body placeholder (idx != 0, i.e. not the title) carries real
    prose — the bullets/agenda archetypes. These slides are text, not tiles; they get
    the relaxed DEADSPACE_TEXT_FORWARD_MAX threshold."""
    for sh in slide.shapes:
        if not getattr(sh, "is_placeholder", False):
            continue
        try:
            idx = sh.placeholder_format.idx
        except Exception:
            continue
        if idx and getattr(sh, "has_text_frame", False) and len(sh.text_frame.text.strip()) >= 80:
            return True
    return False


def check_deadspace(prs):
    """Flag slides whose content leaves a large under-filled content region — see
    module comment above for the fill-ratio model and calibration. Cover slides
    (slide 1) are exempt: a title-only opening slide is a legitimate design."""
    defects = []
    W = in_(prs.slide_width)
    region_area = W * (DEADSPACE_BOTTOM - DEADSPACE_TOP)
    for i, slide in enumerate(prs.slides, 1):
        if i == 1:
            continue
        boxes = _deadspace_bboxes(slide, region_area)
        fill_ratio, empty_frac, rc = _deadspace_grid_metrics(boxes, W)
        limit = DEADSPACE_TEXT_FORWARD_MAX if _is_text_forward(slide) else DEADSPACE_FILL_MAX
        if fill_ratio < limit:
            where = _band_label(rc, DEADSPACE_GRID_W, DEADSPACE_GRID_H) if rc else "diffuse"
            defects.append(Defect(i, "deadspace",
                f"content fills only {fill_ratio:.0%} of the slide; largest empty region "
                f"is {where} ({empty_frac:.0%} of content area)", None))
    return defects


# ---- logo -----------------------------------------------------------------
# Calibrated against deck.py's actual output (see singulator-15b.pptx): the cover and
# section-header layouts bake their own logo picture directly into the layout XML; every
# other light layout inherits the master's logo picture; only the drawn dark "Blank"
# canvases (study_findings, compose(dark=True)) re-embed a logo picture per-slide because
# their full-bleed navy rectangle would otherwise cover the master's inherited one. So
# every slide — cover included — legitimately carries a logo; default policy is "all".

def _sha1_bytes(b):
    return hashlib.sha1(b).hexdigest()


def _template_media_hashes(template_path):
    z = zipfile.ZipFile(template_path)
    hashes = {_sha1_bytes(z.read(n)) for n in z.namelist()
              if n.startswith("ppt/media/") and n.lower().endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff"))}
    z.close()
    return hashes


def _find_logo_shape(shapes, allowed_hashes):
    for sh in iter_shapes(shapes):
        if _stype(sh) == MSO_SHAPE_TYPE.PICTURE:
            try:
                h = _sha1_bytes(sh.image.blob)
            except Exception:
                continue
            if h in allowed_hashes:
                return sh
    return None


def check_logo(prs, template_path, logo_path, require_on):
    if not logo_path:
        return []
    allowed = {_sha1_bytes(open(logo_path, "rb").read())} | _template_media_hashes(template_path)
    defects, direct_positions = [], []
    for i, slide in enumerate(prs.slides, 1):
        direct = _find_logo_shape(slide.shapes, allowed)
        found = direct is not None
        if not found:
            found = (_find_logo_shape(slide.slide_layout.shapes, allowed) is not None
                      or _find_logo_shape(slide.slide_layout.slide_master.shapes, allowed) is not None)
        needs = {"all": True, "noncover": i != 1, "cover": i == 1, "none": False}[require_on]
        if needs and not found:
            defects.append(Defect(i, "logo", "no brand logo found (direct or inherited from layout/master)", None))
        if direct is not None:
            direct_positions.append((i, direct.name, in_(direct.left), in_(direct.top)))
    # position consistency: only among per-slide EXPLICIT logo pictures — layout/master
    # inherited logos are fixed by the template design and legitimately differ by layout.
    if direct_positions:
        bx, by = direct_positions[0][2], direct_positions[0][3]
        for si, name, x, y in direct_positions:
            if abs(x - bx) > LOGO_POS_TOL or abs(y - by) > LOGO_POS_TOL:
                defects.append(Defect(si, "logo",
                    f"logo position ({x:.2f},{y:.2f})in differs from baseline ({bx:.2f},{by:.2f})in", name))
    return defects


# ---- orchestration --------------------------------------------------------

def lint(deck_path, template_path, config_path=None, logo_path=None, require_logo_on="all",
         check_deadspace_flag=True):
    prs = Presentation(deck_path)
    colors, fonts = build_allowed(template_path, config_path)
    defects = []
    defects += check_canvas_size(prs)
    defects += check_palette(prs, colors)
    defects += check_fonts(prs, fonts)
    defects += check_footer(prs)
    defects += check_canvas(prs)
    defects += check_collision(prs)
    defects += check_overflow(prs)
    if check_deadspace_flag:
        defects += check_deadspace(prs)
    defects += check_logo(prs, template_path, logo_path, require_logo_on)
    return prs, defects


# ---- selfcheck --------------------------------------------------------------

def _build_broken_deck(template_path, out_path):
    """Deliberately-broken 2-slide deck: slide 1 is a clean cover (must NOT trigger
    footer/page-number defects); slide 2 bundles one violation of each of the other
    six checks."""
    prs = Presentation(template_path)
    lst = prs.slides._sldIdLst
    for sld in list(lst):
        lst.remove(sld)
    layouts = {L.name: L for L in prs.slide_layouts}

    s1 = prs.slides.add_slide(layouts.get("Title Slide", prs.slide_layouts[0]))
    for ph in s1.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text_frame.text = "Selfcheck Cover"

    s2 = prs.slides.add_slide(layouts.get("Blank", prs.slide_layouts[-1]))

    bad = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(I(0.5)), Emu(I(0.5)), Emu(I(1.0)), Emu(I(1.0)))
    bad.fill.solid(); bad.fill.fore_color.rgb = RGBColor.from_string("FF00AA")   # palette defect
    bad.line.fill.background()

    tb = s2.shapes.add_textbox(Emu(I(2.0)), Emu(I(0.5)), Emu(I(3.0)), Emu(I(0.5)))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = "Off-brand font run"; r.font.name = "Comic Sans MS"; r.font.size = Pt(14)   # fonts defect
    # NOTE: intentionally no footer / page-number textbox on this slide -> footer defect

    s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(I(14.0)), Emu(I(0.5)), Emu(I(1.0)), Emu(I(0.5)))  # canvas defect

    ov = s2.shapes.add_textbox(Emu(I(0.5)), Emu(I(2.0)), Emu(I(3.0)), Emu(I(1.0)))   # overflow defect
    otf = ov.text_frame; otf.word_wrap = True
    for k in range(20):
        p = otf.paragraphs[0] if k == 0 else otf.add_paragraph()
        rr = p.add_run(); rr.text = f"Long filler line forcing wraps and overflow, number {k}."; rr.font.size = Pt(14)

    c1 = s2.shapes.add_textbox(Emu(I(6.0)), Emu(I(2.0)), Emu(I(2.0)), Emu(I(1.0)))   # collision defect
    c1.text_frame.text = "Overlap box A with a real footprint"
    c2 = s2.shapes.add_textbox(Emu(I(7.0)), Emu(I(2.4)), Emu(I(2.0)), Emu(I(1.0)))
    c2.text_frame.text = "Overlap box B with a real footprint"

    prs.save(out_path)
    return out_path


def _build_deadspace_decks(template_path):
    """In-memory (unsaved) deck for check_deadspace directly: slide 1 a dummy cover
    (exempt regardless of content), slide 2 sparse (one small text box top-left of an
    otherwise empty canvas -> must flag), slide 3 tiled with content (colored cards
    covering most of the content band -> must NOT flag)."""
    prs = Presentation(template_path)
    lst = prs.slides._sldIdLst
    for sld in list(lst):
        lst.remove(sld)
    layouts = {L.name: L for L in prs.slide_layouts}
    blank = layouts.get("Blank", prs.slide_layouts[-1])

    prs.slides.add_slide(layouts.get("Title Slide", prs.slide_layouts[0]))

    sparse = prs.slides.add_slide(blank)
    tb = sparse.shapes.add_textbox(Emu(I(0.6)), Emu(I(0.6)), Emu(I(1.5)), Emu(I(0.4)))
    tb.text_frame.text = "Sparse"

    def _tile(slide, with_text):
        cols, rows = 4, 2
        cw, ch = CANVAS_W / cols, (DEADSPACE_BOTTOM - DEADSPACE_TOP) / rows
        for r in range(rows):
            for c in range(cols):
                x, y = c * cw + 0.05, DEADSPACE_TOP + r * ch + 0.05
                box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(I(x)), Emu(I(y)),
                                             Emu(I(cw - 0.1)), Emu(I(ch - 0.1)))
                box.fill.solid(); box.fill.fore_color.rgb = RGBColor.from_string("E8EEF6")
                box.line.fill.background()
                if with_text:   # a separate text box sitting on the tile, like deck.py's stat
                    t = slide.shapes.add_textbox(Emu(I(x + 0.1)), Emu(I(y + 0.1)),
                                                 Emu(I(cw - 0.3)), Emu(I(ch - 0.3)))
                    t.text_frame.text = "Real content on this tile that fills it"

    # slide 3: tiles WITH text on them -> legitimately filled, must NOT flag.
    _tile(prs.slides.add_slide(blank), with_text=True)
    # slide 4: identical tiles but EMPTY (no text) -> a blank slide dressed in filler
    # boxes; the check must NOT be gamed by them (Codex finding, 2026-07-03) -> must flag.
    _tile(prs.slides.add_slide(blank), with_text=False)

    # slide 5: two tall filled cards whose text lives INSIDE the same shape (how
    # deck.py compose `card` blocks are built) -> the card's full box is deliberate
    # ink and must earn full credit, not just its top text band (field report #9).
    cards = prs.slides.add_slide(blank)
    for k in range(2):
        box = cards.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Emu(I(0.6 + k * 6.4)), Emu(I(DEADSPACE_TOP + 0.2)),
                                     Emu(I(5.9)), Emu(I(4.0)))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor.from_string("E8EEF6")
        box.line.fill.background()
        box.text_frame.text = "CARD HEADING"

    return prs


def run_selfcheck():
    import tempfile
    template = DEFAULT_PCS_TEMPLATE
    if not os.path.exists(template):
        print(f"selfcheck: template not found: {template}", file=sys.stderr)
        return False
    fd, path = tempfile.mkstemp(suffix="_lint_selfcheck.pptx")
    os.close(fd)
    try:
        _build_broken_deck(template, path)
        _, defects = lint(path, template)
        found = {d.check for d in defects}
        expect = {"palette", "fonts", "footer", "canvas", "collision", "overflow"}
        all_ok = True
        for chk in sorted(expect):
            ok = chk in found
            all_ok &= ok
            print(f"  [{'PASS' if ok else 'FAIL'}] {chk} defect detected")
        # cover slide must stay clean on the footer check specifically
        cover_footer_clean = not any(d.check == "footer" and d.slide == 1 for d in defects)
        print(f"  [{'PASS' if cover_footer_clean else 'FAIL'}] cover slide has no footer/page-number defect")
        all_ok &= cover_footer_clean

        prs_ds = _build_deadspace_decks(template)
        ds_defects = check_deadspace(prs_ds)
        sparse_flagged = any(d.slide == 2 for d in ds_defects)
        tiled_quiet = not any(d.slide == 3 for d in ds_defects)
        empty_tiles_flagged = any(d.slide == 4 for d in ds_defects)
        cards_quiet = not any(d.slide == 5 for d in ds_defects)
        print(f"  [{'PASS' if sparse_flagged else 'FAIL'}] deadspace flags a sparse top-left-only slide")
        all_ok &= sparse_flagged
        print(f"  [{'PASS' if tiled_quiet else 'FAIL'}] deadspace stays quiet on a tiled slide with real text")
        all_ok &= tiled_quiet
        print(f"  [{'PASS' if empty_tiles_flagged else 'FAIL'}] deadspace flags empty filler tiles (not gamed)")
        all_ok &= empty_tiles_flagged
        print(f"  [{'PASS' if cards_quiet else 'FAIL'}] deadspace credits filled text-bearing cards at full box")
        all_ok &= cards_quiet
        return bool(all_ok)
    finally:
        os.unlink(path)


# ---- CLI --------------------------------------------------------------------

def main(argv=None):
    ap = argparse.ArgumentParser(description="Brand linter for decks built on the deck.py template engine.")
    ap.add_argument("deck", nargs="?")
    ap.add_argument("--template")
    ap.add_argument("--config")
    ap.add_argument("--logo")
    ap.add_argument("--require-logo-on", choices=["all", "noncover", "cover", "none"], default="all")
    ap.add_argument("--no-deadspace", action="store_true", help="skip the dead-space (under-filled slide) check")
    ap.add_argument("--json", action="store_true")
    ap.add_argument("--selfcheck", action="store_true")
    args = ap.parse_args(argv)

    if args.selfcheck:
        print("lint_brand.py --selfcheck")
        ok = run_selfcheck()
        print("ALL PASS" if ok else "SOME FAILED")
        return 0 if ok else 1

    if not args.deck or not args.template:
        print("usage: lint_brand.py DECK.pptx --template TEMPLATE.pptx [--config ...] [--logo ...] "
              "[--require-logo-on all|noncover|cover|none] [--no-deadspace] [--json] [--selfcheck]", file=sys.stderr)
        return 2
    if not os.path.exists(args.deck):
        print(f"error: deck not found: {args.deck}", file=sys.stderr)
        return 2
    if not os.path.exists(args.template):
        print(f"error: template not found: {args.template}", file=sys.stderr)
        return 2

    try:
        prs, defects = lint(args.deck, args.template, args.config, args.logo, args.require_logo_on,
                             check_deadspace_flag=not args.no_deadspace)
    except Exception as e:
        print(f"error: {e}", file=sys.stderr)
        return 2

    if args.json:
        out = {"pass": not defects, "deck": args.deck, "slides": len(prs.slides),
               "defects": [d._asdict() for d in defects]}
        print(json.dumps(out, indent=2))
    else:
        print(f"{args.deck}: {len(prs.slides)} slides, {len(defects)} defect(s)")
        for d in defects:
            loc = f"slide {d.slide}" if d.slide else "deck"
            shape = f" (shape: {d.shape})" if d.shape else ""
            print(f"  {loc}: [{d.check}] {d.detail}{shape}")
        if not defects:
            print("  PASS — no brand defects found.")
    return 0 if not defects else 1


if __name__ == "__main__":
    sys.exit(main())
