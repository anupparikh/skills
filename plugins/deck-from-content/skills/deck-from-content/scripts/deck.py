#!/usr/bin/env python3
"""deck.py — build an on-brand PowerPoint deck from CONTENT ONLY.

A calling agent supplies text and image paths through semantic methods (cover, section,
bullets, comparison, figure, stat_cards, steps, ...). It never touches a layout, colour,
font, or coordinate — those live in the brand template + this engine, so content cannot go
off-brand. Every method that maps to a native template layout fills that layout's
placeholders (so titles, bullets, logo, and background come from the template for free);
only archetypes with no native equivalent (stat cards, numbered steps) are drawn, and those
read the template's own theme colours so they stay on-brand too.

  from deck import Deck
  d = Deck(brand_config="clients/pcs/brand/brand_config.json")   # or template="...pptx"
  d.cover("Singulator", "Automated tissue dissociation")
  d.bullets("What labs need", ["Reproducible", "Gentle", "Automated"])
  d.figure("Nuclei yield", image="fig1.png", caption="Intact nuclei, snRNA-seq ready")
  d.save("out.pptx")
"""
import os, re, json, zipfile, tempfile
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

EMU_IN = 914400


def _hex(v):
    return RGBColor.from_string(str(v).lstrip("#").upper())


def I(v):
    """inches -> EMU int"""
    return int(v * EMU_IN)


class Theme:
    """The 6 accent + neutral colours and 2 fonts, read from the template's theme."""
    def __init__(self, template_path):
        z = zipfile.ZipFile(template_path)
        x = z.read("ppt/theme/theme1.xml").decode()
        z.close()
        def slot(name):
            m = re.search(rf'<a:{name}><a:srgbClr val="([0-9A-Fa-f]{{6}})"/>', x)
            return m.group(1) if m else None
        self.navy   = slot("dk2") or "001A4D"
        self.white  = slot("lt1") or "FFFFFF"
        self.slate  = slot("dk1") or "4A5568"
        self.tint   = slot("lt2") or "E8EEF6"
        self.primary= slot("accent1") or "0056B7"
        self.accent = slot("accent2") or "00A3FF"
        self.green  = slot("accent3") or "26A74A"
        self.red    = slot("accent4") or "DC3445"
        self.card   = slot("accent5") or "16335E"
        mj = re.search(r'<a:majorFont><a:latin typeface="([^"]*)"', x)
        mn = re.search(r'<a:minorFont><a:latin typeface="([^"]*)"', x)
        self.display = mj.group(1) if mj else "Merriweather"
        self.body = mn.group(1) if mn else "Franklin Gothic Book"


# Per-brand, non-paint knobs (see references/brand-manifest.md). These are the values
# deck.py hardcoded before D-034 — kept here as the fallback so a caller that never
# passes manifest= gets byte-identical output to before this existed.
DEFAULT_MANIFEST = {
    "dark_canvas_layout": "Blank",
    "logo": {"dark_position": "top-right", "dark_size_in": 0.38, "keepout_in": 1.2},
    "footer_text": None,
    "title_sizes": {"problem_flow": 30, "study_intro": 25, "dark": 22, "statement": 32, "compose": 24},
    "colors": {
        "muted": "#C8D8EC",
        "page_number": "#9AA5B1",
        "quote_on_light": "#4A5568",
        "quote_on_dark": "#C8D8EC",
        "warn": "#FD7E14",
        "warn_tint": "#FBEAEC",
    },
}


def _load_manifest(manifest):
    """manifest: None | path | dict. Shallow-merges onto DEFAULT_MANIFEST; the three
    nested dicts (logo, title_sizes, colors) merge key-by-key so a brand manifest only
    needs to override what it changes."""
    if manifest is None:
        m = {}
    elif isinstance(manifest, dict):
        m = manifest
    else:
        m = json.load(open(manifest))
    out = {**DEFAULT_MANIFEST, **m}
    for k in ("logo", "title_sizes", "colors"):
        out[k] = {**DEFAULT_MANIFEST[k], **m.get(k, {})}
    return out


def _default_pcs_template():
    return os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "..", "..", "clients", "pcs", "brand", "precision-cell-systems-deck-template.pptx")


def _ensure_template(brand_config, template, out_dir):
    """Resolve a template path. If only brand_config given, generate one via
    brand-deck-template so the caller never has to know templates exist."""
    if template and os.path.exists(template):
        return template
    if brand_config:
        import sys
        gen_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(
            os.path.abspath(__file__)))), "brand-deck-template", "scripts")
        sys.path.insert(0, gen_dir)
        import build_template as B
        r = B.generate(brand_config, out_dir or os.path.dirname(brand_config), None, False)
        return r["pptx"]
    d = _default_pcs_template()
    if os.path.exists(d):
        return d
    raise FileNotFoundError("no template: pass template= or brand_config=")


class Deck:
    def __init__(self, brand_config=None, template=None, footer=None, work_dir=None, manifest=None):
        self.template = _ensure_template(brand_config, template, work_dir)
        self.prs = Presentation(self.template)
        self._clear_slides()
        self.theme = Theme(self.template)

        self.manifest = _load_manifest(manifest)
        # Color resolution order per key: explicit manifest value > (bare Deck: PCS default,
        # keeps legacy output byte-identical) > theme-derived. NEVER inject a PCS hex into a
        # non-PCS brand that merely omitted the key — that leaked PCS colors (Codex finding,
        # 2026-07-03). A brand's own theme colors are always on-palette.
        _raw = ({} if manifest is None else manifest if isinstance(manifest, dict)
                else json.load(open(manifest))).get("colors", {})

        def _color(key, theme_fallback):
            v = _raw.get(key)
            if v is None and manifest is None:
                v = DEFAULT_MANIFEST["colors"][key]        # bare Deck() == PCS, unchanged
            return v.lstrip("#").upper() if v else theme_fallback.lstrip("#").upper()

        self.muted = _color("muted", self.theme.tint)
        self.page_number = _color("page_number", self.theme.slate)
        self.quote_on_light = _color("quote_on_light", self.theme.slate)
        self.quote_on_dark = _color("quote_on_dark", self.theme.tint)
        self.warn = _color("warn", self.theme.accent)
        # pale caution tint for comparison warn footers; theme-derived fallback (a
        # near-white blend of the brand red) so a manifest that omits it stays on-brand.
        self.warn_tint = _color("warn_tint", self._tint_of(self.theme.red, 0.9))
        self.title_sizes = self.manifest["title_sizes"]
        self.dark_canvas_layout = self.manifest["dark_canvas_layout"]
        self.logo_dark_position = self.manifest["logo"]["dark_position"]
        self.logo_dark_size_in = self.manifest["logo"]["dark_size_in"]
        self.logo_keepout_in = self.manifest["logo"]["keepout_in"]
        # explicit footer= kwarg wins (today's behavior); else fall back to the manifest's
        # footer_text (None by default, so a bare Deck() still draws no footer text).
        self.footer = footer if footer is not None else self.manifest.get("footer_text")

        self.W = self.prs.slide_width
        self.H = self.prs.slide_height
        self._layouts = {L.name: L for L in self.prs.slide_layouts}
        self._rev_logo = self._extract_media("image2.png")   # reversed/white logo for dark canvases

    def _extract_media(self, name):
        try:
            z = zipfile.ZipFile(self.template)
            data = z.read(f"ppt/media/{name}"); z.close()
            fd, p = tempfile.mkstemp(suffix="_" + name)
            os.write(fd, data); os.close(fd)
            return p
        except Exception:
            return None

    # ---- infrastructure -------------------------------------------------
    def _clear_slides(self):
        lst = self.prs.slides._sldIdLst
        for s in list(lst):
            lst.remove(s)

    def _add(self, layout_name):
        return self.prs.slides.add_slide(self._layouts[layout_name])

    def _new_slide(self, layout_name, footer=True, dark=False):
        """The one place native-layout slides get created. Stamps footer+page-number
        here so every archetype gets it for free instead of remembering to call
        _footer() itself (root cause of section() and others silently dropping it).
        Dark drawn canvases (_dark_base) draw their own full-bleed background AFTER
        creation, so they stamp footer themselves at the end to keep z-order correct."""
        s = self._add(layout_name)
        if footer:
            self._footer(s, dark=dark)
        return s

    def _ph(self, slide, idx):
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == idx:
                return ph
        return None

    def _set_text(self, ph, text, bold=False):
        tf = ph.text_frame
        tf.text = text if text is not None else ""
        if bold:
            for p in tf.paragraphs:
                for r in p.runs:
                    r.font.bold = True

    def _set_bullets(self, ph, items):
        """items: list of str, or (text, level) tuples, or {'text','level','bold'} dicts."""
        tf = ph.text_frame
        tf.clear()
        first = True
        for it in items:
            if isinstance(it, dict):
                text, level, bold = it.get("text", ""), it.get("level", 0), it.get("bold", False)
            elif isinstance(it, (tuple, list)):
                text, level = it[0], (it[1] if len(it) > 1 else 0)
                bold = False
            else:
                text, level, bold = it, 0, False
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = text
            p.level = level
            if bold:
                for r in p.runs:
                    r.font.bold = True

    def _place_image(self, slide, path, x, y, w, h):
        """Insert a picture scaled to FIT (x,y,w,h) preserving aspect, centred."""
        from PIL import Image
        iw, ih = Image.open(path).size
        box_ar, img_ar = w / h, iw / ih
        if img_ar > box_ar:
            nw = w; nh = int(w / img_ar)
        else:
            nh = h; nw = int(h * img_ar)
        nx = x + (w - nw) // 2
        ny = y + (h - nh) // 2
        return slide.shapes.add_picture(path, nx, ny, nw, nh)

    def _footer(self, slide, dark=False):
        """Draw footer text + page number directly on every slide. The template's
        footer/slide-number placeholders don't render reliably (LibreOffice/export),
        so draw them like the hand-built source deck did — one chrome, every slide."""
        n = len(self.prs.slides._sldIdLst)
        if self.footer:
            tf = self._tb(slide, I(0.5), I(7.04), I(10.6), I(0.3))
            self._run(tf.paragraphs[0], self.footer, 9,
                      self.muted if dark else self.theme.slate)
        pf = self._tb(slide, I(12.2), I(7.04), I(0.83), I(0.3))
        pp = pf.paragraphs[0]; pp.alignment = PP_ALIGN.RIGHT
        self._run(pp, str(n), 9, self.page_number)

    # ---- NATIVE-LAYOUT methods (fill placeholders; fullest brand fidelity) ----
    def section(self, title, kicker=None):
        s = self._new_slide("Section Header")
        self._set_text(self._ph(s, 0), title)
        if kicker is not None:
            self._set_text(self._ph(s, 1), kicker)
        return s

    def bullets(self, title, points):
        s = self._new_slide("Title and Content")
        self._set_text(self._ph(s, 0), title)
        self._set_bullets(self._ph(s, 1), points)
        return s

    def statement(self, title):
        return self._light_title(title, size=self.title_sizes.get("statement", 32))

    def two_column(self, title, left, right):
        s = self._new_slide("Two Content")
        self._set_text(self._ph(s, 0), title)
        self._set_bullets(self._ph(s, 1), left)
        self._set_bullets(self._ph(s, 2), right)
        return s

    def comparison(self, title, a_label, a_points, b_label, b_points):
        s = self._new_slide("Comparison")
        self._set_text(self._ph(s, 0), title)
        self._set_text(self._ph(s, 1), a_label, bold=True)
        self._set_bullets(self._ph(s, 2), a_points)
        self._set_text(self._ph(s, 3), b_label, bold=True)
        self._set_bullets(self._ph(s, 4), b_points)
        return s

    def figure(self, title, image, caption=None):
        """Big image with title + caption beneath — native Picture with Caption."""
        s = self._new_slide("Picture with Caption")
        self._set_text(self._ph(s, 0), title, bold=True)
        pic_ph = self._ph(s, 1)
        try:
            pic_ph.insert_picture(image)
        except Exception:
            self._place_image(s, image, pic_ph.left, pic_ph.top, pic_ph.width, pic_ph.height)
        if caption is not None:
            self._set_text(self._ph(s, 2), caption)
        return s

    def content_with_figure(self, title, image, points):
        """Caption/bullets beside a figure — native Content with Caption."""
        s = self._new_slide("Content with Caption")
        self._set_text(self._ph(s, 0), title, bold=True)
        obj = self._ph(s, 1)
        self._place_image(s, image, obj.left, obj.top, obj.width, obj.height)
        self._set_bullets(self._ph(s, 2), points)
        return s

    # ---- DRAWN on-brand archetypes (no native placeholder equivalent) ----
    def stat_cards(self, title, cards):
        """cards: list of (value, label). Navy tiles with big accent number + label."""
        s = self._new_slide("Title Only")
        self._set_text(self._ph(s, 0), title)
        n = len(cards)
        M = Emu(int(0.5 * EMU_IN)); gap = Emu(int(0.3 * EMU_IN))
        top = Emu(int(2.1 * EMU_IN)); ch = Emu(int(2.6 * EMU_IN))
        avail = self.W - 2 * M - gap * (n - 1)
        cw = int(avail / n)
        x = int(M)
        for value, label in cards:
            self._card(s, x, int(top), cw, int(ch), value, label)
            x += cw + int(gap)
        return s

    def _card(self, slide, x, y, w, h, value, label):
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        box.fill.solid(); box.fill.fore_color.rgb = _hex(self.theme.navy)
        box.line.fill.background()
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = Emu(int(0.25 * EMU_IN)); tf.margin_right = Emu(int(0.25 * EMU_IN))
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = str(value)
        r.font.size = Pt(40); r.font.bold = True
        r.font.color.rgb = _hex(self.theme.accent); r.font.name = self.theme.display
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run(); r2.text = str(label)
        r2.font.size = Pt(12); r2.font.color.rgb = _hex(self.theme.white)
        r2.font.name = self.theme.body

    def steps(self, title, steps):
        """steps: list of (head, desc). Numbered circles + head + description rows."""
        s = self._new_slide("Title Only")
        self._set_text(self._ph(s, 0), title)
        M = int(0.5 * EMU_IN); top = int(2.1 * EMU_IN)
        row_h = int(0.95 * EMU_IN); dia = int(0.5 * EMU_IN)
        y = top
        for i, (head, desc) in enumerate(steps, 1):
            c = s.shapes.add_shape(MSO_SHAPE.OVAL, M, y, dia, dia)
            c.fill.solid(); c.fill.fore_color.rgb = _hex(self.theme.primary)
            c.line.fill.background()
            ctf = c.text_frame; ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
            cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
            cr = cp.add_run(); cr.text = str(i)
            cr.font.size = Pt(20); cr.font.bold = True
            cr.font.color.rgb = _hex(self.theme.white); cr.font.name = self.theme.body
            tx = s.shapes.add_textbox(M + dia + int(0.25 * EMU_IN), y,
                                      self.W - 2 * M - dia - int(0.25 * EMU_IN), row_h)
            ttf = tx.text_frame; ttf.word_wrap = True
            hp = ttf.paragraphs[0]
            hr = hp.add_run(); hr.text = head
            hr.font.size = Pt(15); hr.font.bold = True
            hr.font.color.rgb = _hex(self.theme.navy); hr.font.name = self.theme.body
            if desc:
                dp = ttf.add_paragraph()
                dr = dp.add_run(); dr.text = desc
                dr.font.size = Pt(12); dr.font.color.rgb = _hex(self.theme.slate)
                dr.font.name = self.theme.body
            y += row_h + int(0.15 * EMU_IN)
        return s

    # ==== low-level primitives ===========================================
    def _tb(self, slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
        tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
        tf.word_wrap = wrap; tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = Emu(0); tf.margin_top = tf.margin_bottom = Emu(0)
        return tf

    def _run(self, p, text, size, color, bold=False, italic=False, font=None, spc=None):
        r = p.add_run(); r.text = text; f = r.font
        f.size = Pt(size); f.bold = bold; f.italic = italic
        f.color.rgb = _hex(color); f.name = font or self.theme.body
        if spc is not None:
            r._r.get_or_add_rPr().set("spc", str(spc))
        return r

    def _round(self, slide, x, y, w, h, fill=None, line=None, line_w=1.0, radius=0.08):
        sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
        if fill is None:
            sp.fill.background()
        else:
            sp.fill.solid(); sp.fill.fore_color.rgb = _hex(fill)
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = _hex(line); sp.line.width = Pt(line_w)
        sp.shadow.inherit = False
        return sp

    def _bullets_tb(self, tf, items, color, size=13, gap=6, bullet_color=None):
        bc = bullet_color or self.theme.accent
        for i, it in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(gap)
            self._run(p, "▪  ", size, bc, bold=True)
            self._run(p, it, size, color)

    # ==== canvas bases ====================================================
    def _dark_base(self, title=None):
        """The one place dark full-bleed canvases get created (study_findings,
        compose(dark=True)). Footer is stamped LAST, after the opaque background
        rectangle, so it isn't painted over (an opaque bg drawn after _new_slide's
        immediate footer stamp would hide it — z-order matters here, unlike the light
        native layouts)."""
        s = self._add(self.dark_canvas_layout)
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.W, self.H)
        bg.fill.solid(); bg.fill.fore_color.rgb = _hex(self.theme.navy)
        bg.line.fill.background(); bg.shadow.inherit = False
        if self._rev_logo:
            from PIL import Image
            iw, ih = Image.open(self._rev_logo).size
            h = I(self.logo_dark_size_in); w = int(h * iw / ih)
            lx = I(0.5) if self.logo_dark_position == "top-left" else self.W - I(0.5) - w
            s.shapes.add_picture(self._rev_logo, lx, I(0.4), w, h)
        if title is not None:
            keep = I(self.logo_keepout_in)
            tx, max_tw = I(0.5), I(9.8)
            if self.logo_dark_position == "top-left":
                tx = I(0.5) + keep
                tw = max_tw - keep
            else:
                tw = min(max_tw, self.W - I(0.5) - keep - tx)
            tf = self._tb(s, tx, I(0.34), tw, I(1.15), anchor=MSO_ANCHOR.TOP)
            self._run(tf.paragraphs[0], title, self.title_sizes.get("dark", 22),
                      self.theme.white, bold=True, font=self.theme.display)
        self._footer(s, dark=True)
        return s

    def _light_title(self, title, size=None):
        s = self._new_slide("Title Only")
        ph = self._ph(s, 0)
        self._set_text(ph, title)
        if size:
            for p in ph.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(size)
        return s

    def _no_bullet(self, p):
        pPr = p._p.get_or_add_pPr()
        for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
            for e in pPr.findall(qn(tag)):
                pPr.remove(e)
        pPr.append(pPr.makeelement(qn("a:buNone"), {}))

    def _bullet(self, p, color, char="▪", marL=0.22):
        """Native PowerPoint bullet on a drawn-textbox paragraph: real buChar with a
        hanging indent (marL + negative indent), so wrapped lines align to the text,
        not under the glyph. Call AFTER setting p.space_after (schema child order)."""
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(I(marL)))
        pPr.set("indent", str(-I(marL)))
        for tag in ("a:buNone", "a:buChar", "a:buAutoNum", "a:buClr", "a:buSzPct", "a:buFont"):
            for e in pPr.findall(qn(tag)):
                pPr.remove(e)
        clr = pPr.makeelement(qn("a:buClr"), {})
        clr.append(pPr.makeelement(qn("a:srgbClr"), {"val": str(color).lstrip("#").upper()}))
        pPr.append(clr)
        pPr.append(pPr.makeelement(qn("a:buSzPct"), {"val": "80000"}))
        pPr.append(pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"}))
        pPr.append(pPr.makeelement(qn("a:buChar"), {"char": char}))

    # ==== components (draw on a slide at an explicit box) =================
    def _eyebrow(self, slide, text):
        tf = self._tb(slide, I(0.5), I(0.13), I(9.3), I(0.22))
        self._run(tf.paragraphs[0], text, 10.5, self.theme.primary, bold=True, spc=250)
        return tf

    def _rule(self, slide, x, y, w, color, thick=0.035):
        r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, I(thick))
        r.fill.solid(); r.fill.fore_color.rgb = _hex(color)
        r.line.fill.background(); r.shadow.inherit = False
        return r

    def _chips(self, slide, x, y, labels, prefix=None, max_w=None):
        """max_w: if given, wrap chips onto a new row instead of running past the
        right edge of the cell (compose() passes the grid cell width here)."""
        cx, cy = x, y
        right = (x + max_w) if max_w else None
        if prefix:
            tf = self._tb(slide, cx, cy + I(0.05), I(1.0), I(0.35), anchor=MSO_ANCHOR.MIDDLE)
            self._run(tf.paragraphs[0], prefix, 10, self.theme.slate, bold=True, spc=200)
            cx = x + I(1.05)
        for lab in labels:
            w = I(0.34 + 0.09 * len(lab))
            if right is not None and cx > x and cx + w > right:
                cx, cy = x, cy + I(0.42)
            chip = self._round(slide, cx, cy, w, I(0.34), fill=self.theme.tint, radius=0.5)
            tf = chip.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.word_wrap = False
            tf.margin_top = tf.margin_bottom = Emu(0)
            tf.margin_left = tf.margin_right = Emu(0)
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            self._run(p, lab, 10, self.theme.primary, bold=True)
            cx += w + I(0.12)
        return cx

    def _citation_card(self, slide, x, y, w, h, title, meta=None):
        card = self._round(slide, x, y, w, h, fill=self.theme.tint, radius=0.09)
        tf = card.text_frame; tf.word_wrap = True
        tf.margin_left = I(0.28); tf.margin_right = I(0.28)
        tf.margin_top = I(0.16); tf.margin_bottom = I(0.16)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT          # approved citation cards are left-aligned
        self._run(tf.paragraphs[0], title, 13.5, self.theme.navy, bold=True)
        if meta:
            p = tf.add_paragraph(); p.space_before = Pt(6); p.alignment = PP_ALIGN.LEFT
            self._run(p, meta, 10.5, self.theme.slate)
        return card

    def _figure_panel(self, slide, x, y, w, h, image, caption=None, border=None):
        # border: on a light slide a white panel is invisible against the white
        # canvas — a thin tint edge gives the figure a frame (approved decks do this).
        panel = self._round(slide, x, y, w, h, fill=self.theme.white, line=border,
                            line_w=1.0, radius=0.04)
        pad = I(0.18)
        cap_h = I(0.5) if caption else 0
        self._place_image(slide, image, x + pad, y + pad, w - 2 * pad, h - 2 * pad - cap_h)
        if caption:
            tf = self._tb(slide, x + pad, y + h - cap_h - I(0.02), w - 2 * pad, cap_h)
            self._run(tf.paragraphs[0], caption, 9, self.theme.slate, italic=True)
        return panel

    def _findings(self, slide, x, y, w, h, points, label="KEY FINDINGS"):
        tf = self._tb(slide, x, y, w, h)
        self._run(tf.paragraphs[0], label, 11, self.theme.accent, bold=True, spc=200)
        p = tf.add_paragraph(); p.space_before = Pt(6)
        # first findings paragraph follows the label paragraph
        for i, it in enumerate(points):
            pp = tf.add_paragraph() if i else p
            pp.space_after = Pt(7)
            self._bullet(pp, self.theme.accent)
            self._run(pp, it, 12.5, self.theme.white)
        return tf

    def _big_number(self, slide, x, y, w, value, label=None, color=None, size=54, h=1.4,
                    label_color=None):
        col = color or self.theme.accent
        tf = self._tb(slide, x, y, w, I(h))
        self._run(tf.paragraphs[0], str(value), size, col, bold=True, font=self.theme.display)
        if label:
            p = tf.add_paragraph(); p.space_before = Pt(4)
            self._run(p, label, 12, label_color or self.theme.white)
        return tf

    def _statement(self, slide, x, y, w, text, sub=None, h=2.2):
        tf = self._tb(slide, x, y, w, I(h))
        self._run(tf.paragraphs[0], text, 34, self.theme.white, bold=True, font=self.theme.display)
        if sub:
            p = tf.add_paragraph(); p.space_before = Pt(6)
            self._run(p, sub, 13, self.muted)
        return tf

    def _quote_bar(self, slide, x, y, w, h, text, dark=True):
        """dark: whether this quote sits on a dark or light background — picks
        quote_on_dark (light-on-navy) vs quote_on_light (slate) so the text stays
        readable either way instead of always painting the on-navy color."""
        self._rule(slide, x, y, I(0.06), self.theme.accent, thick=h / EMU_IN)
        tf = self._tb(slide, x + I(0.22), y, w - I(0.22), h, anchor=MSO_ANCHOR.MIDDLE)
        color = self.quote_on_dark if dark else self.quote_on_light
        self._run(tf.paragraphs[0], text, 13, color, italic=True)
        return tf

    @staticmethod
    def _tint_of(hexstr, toward_white=0.9):
        """Blend a theme color toward white (no literal hex in code — a caution/warn
        tint is derived from the brand's own red, not a hardcoded pink)."""
        h = str(hexstr).lstrip("#")
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        f = toward_white
        r = int(r + (255 - r) * f); g = int(g + (255 - g) * f); b = int(b + (255 - b) * f)
        return f"{r:02X}{g:02X}{b:02X}"

    def _callout_bar(self, slide, x, y, w, h, text):
        bar = self._round(slide, x, y, w, h, fill=self.theme.primary, radius=0.12)
        tf = bar.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = I(0.25); tf.margin_right = I(0.25)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        self._run(p, text, 14, self.theme.white, bold=True)
        return bar

    def _process_flow(self, slide, x, y, w, h, cards):
        n = len(cards); arrow = I(0.4)
        cw = int((w - arrow * (n - 1)) / n)
        cx = x
        for i, (head, desc) in enumerate(cards):
            highlight = (i == n - 1)
            box = self._round(slide, cx, y, cw, h,
                              fill=(self.theme.tint if not highlight else self.theme.white),
                              line=(self.theme.primary if highlight else None),
                              line_w=1.5, radius=0.08)
            tf = box.text_frame; tf.word_wrap = True
            tf.margin_left = I(0.12); tf.margin_right = I(0.12)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            self._run(tf.paragraphs[0], head, 12, self.theme.navy, bold=True)
            if desc:
                p = tf.add_paragraph(); p.space_before = Pt(2)
                self._run(p, desc, 9.5, self.theme.slate)
            if i < n - 1:
                atf = self._tb(slide, cx + cw, y, arrow, h, anchor=MSO_ANCHOR.MIDDLE)
                ap = atf.paragraphs[0]; ap.alignment = PP_ALIGN.CENTER
                self._run(ap, "→", 20, self.theme.primary, bold=True)
            cx += cw + arrow
        return slide

    def _vertical_steps(self, slide, x, y, w, h, steps):
        n = len(steps); gap = I(0.14)
        sh = int((h - gap * (n - 1)) / n)
        sy = y
        for i, (head, desc) in enumerate(steps, 1):
            box = self._round(slide, x, sy, w, sh, fill=self.theme.tint, radius=0.12)
            dia = min(sh - I(0.12), I(0.44))
            c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + I(0.12), sy + (sh - dia) // 2, dia, dia)
            c.fill.solid(); c.fill.fore_color.rgb = _hex(self.theme.primary)
            c.line.fill.background(); c.shadow.inherit = False
            ctf = c.text_frame; ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
            cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
            self._run(cp, str(i), 16, self.theme.white, bold=True)
            tf = self._tb(slide, x + I(0.12) + dia + I(0.16), sy, w - dia - I(0.5), sh,
                          anchor=MSO_ANCHOR.MIDDLE)
            self._run(tf.paragraphs[0], head, 12, self.theme.navy, bold=True)
            if desc:
                p = tf.add_paragraph()
                self._run(p, desc, 9.5, self.theme.slate)
            sy += sh + gap
        return slide

    # ==== composite semantic slides (encode the target archetypes) =======
    def cover(self, title, subtitle=None, tagline=None):
        s = self._new_slide("Title Slide", footer=False)   # cover carries no footer/page-number chrome
        self._rule(s, I(0.5), I(2.9), I(0.7), self.theme.accent)
        self._set_text(self._ph(s, 0), title)
        if subtitle is not None:
            self._set_text(self._ph(s, 1), subtitle)
        if tagline is not None:
            tf = self._tb(s, I(0.5), I(6.3), I(11), I(0.5))
            self._run(tf.paragraphs[0], tagline, 13, self.theme.accent, italic=True)
        return s

    def agenda(self, title, groups):
        """groups: list of (label, [ (item, gloss) ]). Native Title and Content body,
        two sub-headed numbered sections."""
        s = self._new_slide("Title and Content")
        self._set_text(self._ph(s, 0), title)
        body = self._ph(s, 1)
        tf = body.text_frame; tf.clear()
        n = 0; first = True
        for label, items in groups:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            self._no_bullet(p)
            p.space_before = Pt(6 if not first else 0); p.space_after = Pt(3)
            self._run(p, label, 11, self.theme.primary, bold=True, spc=150)
            for item, gloss in items:
                n += 1
                q = tf.add_paragraph(); self._no_bullet(q); q.space_after = Pt(2)
                self._run(q, f"{n}.  ", 10.5, self.theme.navy, bold=True)
                self._run(q, item, 10.5, self.theme.navy, bold=True)
                if gloss:
                    self._run(q, f"  {gloss}", 10.5, self.theme.slate)
        return s

    def problem_flow(self, title, body, cards, answer=None, caption=None, eyebrow=None):
        """Title + framing paragraph + horizontal step-flow + colored answer bar."""
        s = self._light_title(title, size=self.title_sizes.get("problem_flow", 30))
        if eyebrow:
            self._eyebrow(s, eyebrow)
        tf = self._tb(s, I(0.5), I(1.68), I(12.33), I(0.9))
        self._run(tf.paragraphs[0], body, 14, self.theme.slate)
        self._process_flow(s, I(0.5), I(2.78), I(12.33), I(1.95), cards)
        yb = 5.05
        if answer:
            self._callout_bar(s, I(0.5), I(yb), I(12.33), I(0.95), answer); yb += 1.2
        if caption:
            cf = self._tb(s, I(0.5), I(yb), I(12.33), I(0.5))
            self._run(cf.paragraphs[0], caption, 11, self.theme.slate, italic=True)
        return s

    def study_intro(self, title, citation, method=None, challenge_title=None,
                    challenge_points=None, eyebrow=None):
        """White study-intro: title, light-blue citation card, METHOD chips, challenge bullets.
        citation: (paper_title, meta) ; method: list of chip labels."""
        s = self._light_title(title, size=self.title_sizes.get("study_intro", 25))
        if eyebrow:
            self._eyebrow(s, eyebrow)
        # distribute over the full frame (citation high, challenge in the lower half)
        self._citation_card(s, I(0.5), I(1.55), I(12.33), I(1.35),
                            citation[0], citation[1] if len(citation) > 1 else None)
        if method:
            self._chips(s, I(0.5), I(3.25), method, prefix="METHOD")
        if challenge_title or challenge_points:
            tf = self._tb(s, I(0.5), I(4.15), I(12.33), I(2.3))
            if challenge_title:
                self._run(tf.paragraphs[0], challenge_title, 14.5, self.theme.navy, bold=True)
                start = tf.add_paragraph(); start.space_before = Pt(4)
            for i, pt in enumerate(challenge_points or []):
                p = start if (challenge_title and i == 0) else tf.add_paragraph()
                p.space_after = Pt(8)
                self._bullet(p, self.theme.primary)
                self._run(p, pt, 13.5, self.theme.slate)
        return s

    def study_findings(self, title, *, figure=None, figure_caption=None, steps=None,
                       stat=None, stat_label=None, statement=None, statement_sub=None,
                       findings=None, findings_label="KEY FINDINGS", quote=None, why=None,
                       light=False, eyebrow=None):
        """Findings slide: LEFT = figure panel or vertical steps; RIGHT = big
        stat / statement + KEY FINDINGS (+ optional quote); italic why footer.

        Dark canvas by default. light=True renders the approved case-study findings
        archetype: white slide, serif navy headline, figure/steps left, and the
        stat + KEY FINDINGS grouped inside a navy accent card on the right so the
        big number and findings read against a panel (not floating on white)."""
        if figure is not None and not os.path.isabs(figure) and getattr(self, "_img_dir", None):
            figure = os.path.join(self._img_dir, figure)
        if light:
            s = self._light_title(title, size=self.title_sizes.get("study_intro", 25))
            if eyebrow:
                self._eyebrow(s, eyebrow)
        else:
            s = self._dark_base(title)
        LX, LY, LW, LH = I(0.5), I(1.6), I(6.2), I(4.7)
        RX, RY, RW = I(7.0), I(1.6), I(5.83)
        if figure is not None:
            self._figure_panel(s, LX, LY, LW, LH, figure, figure_caption,
                              border=(self.muted if light else None))
        elif steps is not None:
            self._vertical_steps(s, LX, LY, LW, LH, steps)
        # Right column. On light, wrap it in a navy accent card so the white/accent
        # text reads; the card's inner padding defines the working box.
        RH = I(4.7)
        if light:
            self._round(s, RX, RY, RW, RH, fill=self.theme.navy, radius=0.06)
            pad = I(0.34)
            cx, cw, ry, r_bottom = RX + pad, RW - 2 * pad, RY + pad, RY + RH - pad
        else:
            cx, cw, ry, r_bottom = RX, RW, RY, I(6.4)
        if stat is not None:
            # a short token ("143,051") is a hero number; a phrase ("On par") is a
            # headline — size down so it doesn't wrap or crowd the card.
            big = 54 if len(str(stat)) <= 8 else 30
            self._big_number(s, cx, ry, cw, stat, stat_label, size=big, h=1.5); ry += I(1.5)
        elif statement is not None:
            self._statement(s, cx, ry, cw, statement, statement_sub, h=1.9); ry += I(1.9)
        if quote is not None:
            self._quote_bar(s, cx, ry, cw, I(0.9), quote, dark=True); ry += I(1.0)
        if findings:
            self._findings(s, cx, ry, cw, r_bottom - ry, findings, findings_label)
        if why:
            wf = self._tb(s, I(0.5), I(6.55), I(11.5), I(0.4))
            self._run(wf.paragraphs[0], why, 10.5,
                      self.theme.slate if light else self.muted, italic=True)
        return s

    def comparison_columns(self, title, columns, eyebrow=None):
        """Approved "How labs handle X today" archetype: N side-by-side option cards,
        each with a colored header band (label), a method line, a green STRENGTHS
        list of `+` items, and a caution footer (⚠, red-tinted). Light slide.

        columns: list of dicts {label, desc?, strengths?[...], warn?, accent?}. accent
        is a theme key for the header band ('slate'|'primary'|'navy'); defaults walk
        slate→primary→navy so the first (worst) option reads neutral and the last most
        branded — the progression the approved decks use."""
        s = self._light_title(title, size=self.title_sizes.get("compose", 24))
        if eyebrow:
            self._eyebrow(s, eyebrow)
        n = max(len(columns), 1)
        x0, x1, gut = 0.5, 12.83, 0.3
        cw = (x1 - x0 - (n - 1) * gut) / n
        y, ch = 1.7, 4.7
        hdr = 0.5
        accents = {"slate": self.theme.slate, "primary": self.theme.primary, "navy": self.theme.navy}
        default_accents = ["slate", "primary", "navy"]
        warn_fill = self.warn_tint
        for i, col in enumerate(columns):
            cx = x0 + i * (cw + gut)
            # card
            self._round(s, I(cx), I(y), I(cw), I(ch), fill=self.theme.tint,
                        line=self.muted, line_w=0.75, radius=0.04)
            # header band
            acc = accents.get(col.get("accent") or default_accents[min(i, 2)], self.theme.navy)
            band = self._round(s, I(cx), I(y), I(cw), I(hdr), fill=acc, radius=0.04)
            btf = band.text_frame; btf.vertical_anchor = MSO_ANCHOR.MIDDLE
            btf.margin_left = I(0.16)
            btf.paragraphs[0].alignment = PP_ALIGN.LEFT
            self._run(btf.paragraphs[0], col["label"], 12, self.theme.white, bold=True, spc=150)
            iy = y + hdr + 0.2
            # method line
            if col.get("desc"):
                dtf = self._tb(s, I(cx + 0.16), I(iy), I(cw - 0.32), I(0.8))
                self._run(dtf.paragraphs[0], col["desc"], 12.5, self.theme.navy, bold=True)
                iy += 0.95
            # strengths
            strengths = col.get("strengths") or []
            if strengths:
                ltf = self._tb(s, I(cx + 0.16), I(iy), I(cw - 0.32), I(0.24))
                self._run(ltf.paragraphs[0], "STRENGTHS", 10, self.theme.green, bold=True, spc=150)
                iy += 0.3
                stf = self._tb(s, I(cx + 0.16), I(iy), I(cw - 0.32), I(1.4))
                for j, it in enumerate(strengths):
                    p = stf.paragraphs[0] if j == 0 else stf.add_paragraph()
                    p.space_after = Pt(5)
                    self._run(p, f"+  {it}", 11.5, self.theme.slate)
            # caution footer
            if col.get("warn"):
                wh = 1.05
                wy = y + ch - wh - 0.12
                self._round(s, I(cx + 0.14), I(wy), I(cw - 0.28), I(wh), fill=warn_fill, radius=0.06)
                wtf = self._tb(s, I(cx + 0.28), I(wy + 0.08), I(cw - 0.5), I(wh - 0.16))
                self._run(wtf.paragraphs[0], f"⚠  {col['warn']}", 10.5, self.theme.slate)
        return s

    def checklist(self, title, items, eyebrow=None, intro=None):
        """Approved "What single-cell labs need" archetype: a green-check requirement
        list on white. items: list of strings. Rows tile the frame so a 4-8 item list
        fills vertically."""
        s = self._light_title(title, size=self.title_sizes.get("compose", 24))
        if eyebrow:
            self._eyebrow(s, eyebrow)
        y0 = 1.7
        if intro:
            itf = self._tb(s, I(0.5), I(1.62), I(12.0), I(0.5))
            self._run(itf.paragraphs[0], intro, 13, self.theme.slate)
            y0 = 2.25
        n = max(len(items), 1)
        stride = min(0.85, (6.35 - y0) / n)
        box = min(0.36, stride - 0.16)
        for i, it in enumerate(items):
            ry = y0 + i * stride
            chk = self._round(s, I(0.55), I(ry), I(box), I(box), fill=self.theme.green, radius=0.22)
            ctf = chk.text_frame; ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
            ctf.margin_left = ctf.margin_right = ctf.margin_top = ctf.margin_bottom = Emu(0)
            cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
            self._run(cp, "✓", 13, self.theme.white, bold=True)
            ttf = self._tb(s, I(0.55 + box + 0.28), I(ry - 0.04), I(12.83 - (0.55 + box + 0.28)),
                           I(stride), anchor=MSO_ANCHOR.MIDDLE)
            self._run(ttf.paragraphs[0], it, 13.5, self.theme.navy)
        return s

    def dueling_numbers(self, title, left, right, footnote=None, eyebrow=None):
        """Approved "1 vs 8" archetype: two side-by-side cards, each a colored header
        band + a huge number + a caption + a note; a footnote spans below. left/right:
        dict {label, value, caption?, note?, accent?}. Left reads as the win (primary),
        right as the cost (slate) unless accents are given."""
        s = self._light_title(title, size=self.title_sizes.get("compose", 24))
        if eyebrow:
            self._eyebrow(s, eyebrow)
        y, ch, hdr = 1.72, 4.35, 0.5
        cw, gut = 5.9, 0.5
        x0 = 0.5
        accents = {"primary": self.theme.primary, "slate": self.theme.slate, "navy": self.theme.navy}
        for i, (col, cx, defacc) in enumerate([(left, x0, "primary"),
                                               (right, x0 + cw + gut, "slate")]):
            self._round(s, I(cx), I(y), I(cw), I(ch), fill=self.theme.tint, line=self.muted,
                        line_w=0.75, radius=0.04)
            acc = accents.get(col.get("accent") or defacc, self.theme.slate)
            band = self._round(s, I(cx), I(y), I(cw), I(hdr), fill=acc, radius=0.04)
            btf = band.text_frame; btf.vertical_anchor = MSO_ANCHOR.MIDDLE; btf.margin_left = I(0.18)
            btf.paragraphs[0].alignment = PP_ALIGN.LEFT
            self._run(btf.paragraphs[0], col["label"], 12, self.theme.white, bold=True, spc=150)
            ntf = self._tb(s, I(cx + 0.2), I(y + 0.75), I(cw - 0.4), I(1.7), anchor=MSO_ANCHOR.MIDDLE)
            np = ntf.paragraphs[0]; np.alignment = PP_ALIGN.CENTER
            self._run(np, str(col["value"]), 80, acc if acc != self.theme.slate else self.theme.navy,
                      bold=True, font=self.theme.display)
            iy = y + 2.55
            if col.get("caption"):
                ctf = self._tb(s, I(cx + 0.2), I(iy), I(cw - 0.4), I(0.3))
                cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
                self._run(cp, col["caption"], 12.5, self.theme.slate, bold=True)
                iy += 0.42
            if col.get("note"):
                ttf = self._tb(s, I(cx + 0.25), I(iy), I(cw - 0.5), I(ch - (iy - y) - 0.15))
                tp = ttf.paragraphs[0]; tp.alignment = PP_ALIGN.CENTER
                self._run(tp, col["note"], 11.5, self.theme.slate)
        if footnote:
            ftf = self._tb(s, I(0.5), I(y + ch + 0.15), I(12.33), I(0.4))
            fp = ftf.paragraphs[0]; fp.alignment = PP_ALIGN.CENTER
            self._run(fp, footnote, 12, self.theme.navy, bold=True)
        return s

    def closing(self, title="Thank you.", contact=None, dark=True):
        """Approved closing archetype: a large centered serif line + an optional
        centered contact card. Both approved decks render this on dark navy, so dark
        is the default. contact: {name?, detail?, website?}."""
        if dark:
            s = self._dark_base(None)                    # navy bg + logo, no corner title
            title_c, card_fill = self.theme.white, self.theme.card
            name_c, detail_c = self.theme.white, self.muted
        else:
            s = self._new_slide("Title Only", footer=False)
            try:
                sp = self._ph(s, 0)._element; sp.getparent().remove(sp)
            except Exception:
                pass
            title_c, card_fill = self.theme.navy, self.theme.tint
            name_c, detail_c = self.theme.navy, self.theme.slate
        ttf = self._tb(s, I(0.5), I(1.6), I(12.33), I(1.2), anchor=MSO_ANCHOR.MIDDLE)
        tp = ttf.paragraphs[0]; tp.alignment = PP_ALIGN.CENTER
        self._run(tp, title, 46, title_c, bold=True, font=self.theme.display)
        if contact:
            cw, cardh = 5.4, 1.9
            cx = (13.333 - cw) / 2
            card = self._round(s, I(cx), I(3.5), I(cw), I(cardh), fill=card_fill, radius=0.06)
            tf = card.text_frame; tf.word_wrap = True
            tf.margin_left = tf.margin_right = I(0.3); tf.margin_top = I(0.22)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
            self._run(p0, "CONTACT", 10.5, self.theme.accent, bold=True, spc=200)
            if contact.get("name"):
                p = tf.add_paragraph(); p.space_before = Pt(6); p.alignment = PP_ALIGN.CENTER
                self._run(p, contact["name"], 15, name_c, bold=True)
            if contact.get("detail"):
                p = tf.add_paragraph(); p.space_before = Pt(3); p.alignment = PP_ALIGN.CENTER
                self._run(p, contact["detail"], 12, detail_c)
            if contact.get("website"):
                p = tf.add_paragraph(); p.space_before = Pt(4); p.alignment = PP_ALIGN.CENTER
                self._run(p, contact["website"], 12, self.theme.accent)
        return s

    # ==== FREEFORM COMPOSITION (coordinate-free 12-col grid) ==============
    # compose() lets a caller lay out a page creatively without touching a colour,
    # font or coordinate: it places on-brand blocks into a 12-column grid and the
    # engine owns margins, gutters, borders, logo and footer — so any arrangement
    # still comes out on-brand and clean. Complements the fixed semantic methods.
    def compose(self, blocks, title=None, dark=False, eyebrow=None, why=None, images_dir=None):
        """blocks: list of dicts, each {type, col, colspan, row, rowspan, ...fields}.
        type in: text | bullets | card | stat | chips | quote | callout | figure."""
        self._img_dir = images_dir
        if dark:
            s = self._dark_base(title)
        else:
            s = self._light_title(title, size=self.title_sizes.get("compose", 24)) if title \
                else self._new_slide("Title Only")
        if eyebrow:
            self._eyebrow(s, eyebrow)
        ax0, ax1, ay1 = 0.5, 12.83, 6.4
        if title:
            ay0 = 1.62
        elif dark and self._rev_logo:
            # no title means _dark_base drew no title box to keep the grid clear of the
            # logo corner (see _dark_base) — but the logo itself is still there, so start
            # the grid below it rather than under the fixed 0.55in used when there's no logo.
            ay0 = 0.4 + self.logo_dark_size_in + 0.12
        else:
            ay0 = 0.55
        gut, ncols = 0.22, 12
        nrows = max((b.get("row", 0) + b.get("rowspan", 1) for b in blocks), default=1)
        colw = (ax1 - ax0 - (ncols - 1) * gut) / ncols
        rowh = (ay1 - ay0 - (nrows - 1) * gut) / max(1, nrows)
        for b in blocks:
            c, cs = b.get("col", 0), b.get("colspan", 12)
            r, rs = b.get("row", 0), b.get("rowspan", 1)
            x = ax0 + c * (colw + gut); w = cs * colw + (cs - 1) * gut
            y = ay0 + r * (rowh + gut); h = rs * rowh + (rs - 1) * gut
            self._draw_block(s, b, I(x), I(y), I(w), I(h), dark)
        if why:
            tf = self._tb(s, I(0.5), I(6.55), I(11.5), I(0.4))
            self._run(tf.paragraphs[0], why, 10.5, self.muted if dark else self.theme.slate, italic=True)
        return s

    def _draw_block(self, s, b, x, y, w, h, dark):
        t = b["type"]
        muted = self.muted
        if t == "text":
            tf = self._tb(s, x, y, w, h)
            if b.get("title"):
                self._run(tf.paragraphs[0], b["title"], 16,
                          self.theme.white if dark else self.theme.navy, bold=True, font=self.theme.display)
                if b.get("body"):
                    p = tf.add_paragraph(); p.space_before = Pt(4)
                    self._run(p, b["body"], 13, muted if dark else self.theme.slate)
            elif b.get("body"):
                self._run(tf.paragraphs[0], b["body"], 13, muted if dark else self.theme.slate)
        elif t == "bullets":
            tf = self._tb(s, x, y, w, h)
            dot = self.theme.accent if dark else self.theme.primary
            idx0 = 0
            if b.get("label"):
                self._run(tf.paragraphs[0], b["label"], 11, dot, bold=True, spc=200)
                idx0 = 1
            for i, it in enumerate(b["items"]):
                p = tf.paragraphs[0] if (idx0 == 0 and i == 0) else tf.add_paragraph()
                p.space_after = Pt(7)
                self._bullet(p, dot)
                self._run(p, it, 12.5, muted if dark else self.theme.slate)
        elif t == "card":
            variant = b.get("variant") or ("navy" if dark else "tint")
            fill = self.theme.card if variant == "navy" else self.theme.tint
            edge_map = {"accent": self.theme.accent, "accent_bright": self.theme.accent,
                        "primary": self.theme.primary, "warn": self.warn,
                        "green": self.theme.green, "red": self.theme.red, "tint": muted}
            edge = edge_map.get(b.get("accent"), self.theme.accent if variant == "navy" else muted)
            th = self.theme.white if variant == "navy" else self.theme.navy
            tb = self.theme.white if variant == "navy" else self.theme.slate
            box = self._round(s, x, y, w, h, fill=fill, line=edge, line_w=1.25, radius=0.06)
            tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
            tf.margin_left = I(0.14); tf.margin_right = I(0.14); tf.margin_top = I(0.1)
            if b.get("heading"):
                self._run(tf.paragraphs[0], b["heading"], 12, th, bold=True)
                p = tf.add_paragraph(); p.space_before = Pt(3)
                self._run(p, b["body"], 12, tb)
            else:
                self._run(tf.paragraphs[0], b["body"], 12, tb)
        elif t == "stat":
            self._round(s, x, y, w, h, fill=self.theme.card, line=self.theme.accent, line_w=1.25, radius=0.08)
            tf = self._tb(s, x + I(0.18), y + I(0.14), w - I(0.36), h - I(0.28))
            self._run(tf.paragraphs[0], str(b["big"]), b.get("size", 40),
                      self.theme.accent, bold=True, font=self.theme.display)
            if b.get("body"):
                p = tf.add_paragraph(); p.space_before = Pt(4)
                self._run(p, b["body"], 12, self.theme.white)
        elif t == "chips":
            self._chips(s, x, y, b["items"], max_w=w)
        elif t == "quote":
            self._quote_bar(s, x, y, w, h, b["text"], dark=dark)
            if b.get("attrib"):
                af = self._tb(s, x + I(0.22), y + h - I(0.32), w - I(0.22), I(0.3))
                self._run(af.paragraphs[0], b["attrib"], 11,
                          self.quote_on_dark if dark else self.quote_on_light)
        elif t == "callout":
            self._callout_bar(s, x, y, w, h, b["text"])
        elif t == "figure":
            img = b["image"]
            if not os.path.isabs(img) and self._img_dir:
                img = os.path.join(self._img_dir, img)
            self._figure_panel(s, x, y, w, h, img, b.get("caption"))
        elif t == "table":
            cols, rows = b["columns"], b["rows"]
            hl = b.get("highlight_last", True)
            n = len(cols)
            if n <= 1:                       # single column: no label/value split, full width
                xs, ws = [x], [w]
            else:
                first_w = int(w * 0.34); restw = int((w - first_w) / (n - 1))
                xs = [x] + [x + first_w + i * restw for i in range(n - 1)]
                ws = [first_w] + [restw] * (n - 1)
            nr = len(rows) + 1; rh = int(h / nr)
            pad = I(0.1)
            if hl:
                self._round(s, xs[-1] - I(0.06), y, ws[-1] + I(0.06), h,
                            fill=self.theme.card, line=self.theme.accent, line_w=1.5, radius=0.05)
            for ci, ct in enumerate(cols):
                col = self.theme.accent if (hl and ci == n - 1) else self.theme.white
                tf = self._tb(s, xs[ci] + pad, y, ws[ci] - pad, rh, anchor=MSO_ANCHOR.MIDDLE)
                self._run(tf.paragraphs[0], ct, 12, col, bold=True)
            for ri, row in enumerate(rows):
                yy = y + rh * (ri + 1)
                for ci, val in enumerate(row):
                    if ci == 0:
                        col, bold, sz = self.theme.white, True, 12
                    elif hl and ci == n - 1:
                        col, bold, sz = self.theme.accent, True, 12.5
                    else:
                        col, bold, sz = muted, False, 12
                    tf = self._tb(s, xs[ci] + pad, yy, ws[ci] - pad, rh, anchor=MSO_ANCHOR.MIDDLE)
                    self._run(tf.paragraphs[0], str(val), sz, col, bold=bold)
        else:
            raise ValueError(f"unknown block type: {t!r}")

    # ---- output ----------------------------------------------------------
    def save(self, path):
        self.prs.save(path)
        return path
